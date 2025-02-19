import os
import re
import json
import logging
import aisuite as ai
from dotenv import load_dotenv
from typing import Tuple, List, Dict, Optional

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE, PP_PARAGRAPH_ALIGNMENT
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

from bs4 import BeautifulSoup
from tenacity import retry, stop_after_attempt, wait_exponential


# Configure logging
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__) 

# Load environment variables from .env file
load_dotenv()




@retry(stop=stop_after_attempt(3), wait=wait_exponential(multiplier=1, min=4, max=15))
def get_completion(messages: list[dict], model="openai:gpt-4o") -> str:
    """ Generate a completion for the given messages and model.
    
    Args:
        messages (list): A list of messages, where each message is a dictionary with the following keys:
            - role: The role of the sender of the message, e.g. "user" or "system".
            - content: The text of the message.
        model (str): The model to use to generate the completion.
    
    Returns:
        str: The generated completion.
    """
    logger.info(f"Getting completion for messages:")
    client = ai.Client()
    client.configure({"openai" : {
  "api_key": os.environ.get("API_KEY"),
}})
    response = None
    model = model
    try:
        logger.info("Trying to get completion for messages")
        response = client.chat.completions.create(
                model=model,
                messages=messages,
                temperature=1.0,
            )
    except Exception as e:
        logger.error(f"Error getting completion for messages.\nException: {e}")
        raise  # Allow @retry to handle the exception
    else:
        logger.info(f"successfully got completion for messages")
        return response.choices[0].message.content
    

def parse_topics(xml_string: str) -> Tuple[str, str, List[Dict[str, List[str]]]]:
    """
    Parse an XML string containing slides into structured data.

    Args:
        xml_string (str): XML string containing topic suggestions.

    Returns:
        Tuple[str, str, List[Dict[str, List[str]]]]: A tuple containing:
            - Title of the slides (str).
            - Subtitle of the slides (str).
            - A list of dictionaries with slide titles and bullet points.

    Raises:
        ValueError: If XML parsing fails or required elements are missing.
    """
    logger.info("Parsing topics from XML.")
    print(xml_string)

    try:
        # Parse the XML string using BeautifulSoup
        match = re.search(r"<PowerPoint.*?>.*?</PowerPoint>", xml_string, re.DOTALL | re.IGNORECASE)
        xml_clean = match.group(0)
        soup = BeautifulSoup(xml_clean, 'xml')
        logger.debug("XML successfully parsed with BeautifulSoup.")
    except Exception as e:
        logger.error(f"Failed to parse XML: {e}")
        raise ValueError("Invalid XML input.") from e

    try:
        # Extract title and subtitle
        title = soup.find(re.compile('Title', re.IGNORECASE)).get_text(strip=True)
        subtitle = soup.find(re.compile('Subtitle', re.IGNORECASE)).get_text(strip=True)
        logger.debug(f"Extracted Title: {title}, Subtitle: {subtitle}")
    except AttributeError as e:
        logger.error("Missing required elements in the XML: Title or Subtitle.")
        raise ValueError("XML is missing required elements: Title or Subtitle.") from e

    # Extract slides data
    slides_data = []
    for slide in soup.find_all(re.compile('Slide', re.IGNORECASE)):
        slide_title = slide.find(re.compile('Title', re.IGNORECASE))
        bullet_points = slide.find_all(re.compile('BulletPoint', re.IGNORECASE))

        if slide_title:
            slide_dict = {
                "title": slide_title.get_text(strip=True),
                "points": [point.get_text(strip=True) for point in bullet_points]
            }
            slides_data.append(slide_dict)
            logger.debug(f"Parsed slide: {slide_dict}")

    logger.info("Successfully parsed topics from XML.")

    return title, subtitle, slides_data


def extract_json(text: str) -> Optional[dict]:
    """
    Extract JSON data from text string.

    Args:
        text (str): Text containing JSON data.

    Returns:
        Optional[tuple]: A tuple containing:
            - Title of the slides (str)
            - Subtitle of the slides (str) 
            - A list of dictionaries with slide data
        Returns None if no valid JSON is found.

    """
    logger.info("start extracting JSON from text.")

    # Use regex to extract the JSON part
    json_match = re.search(r'\{.*\}', text, re.DOTALL)

    if json_match:
        json_str = json_match.group(0)  # Extract matched JSON string
        data_dict = json.loads(json_str)  # Convert JSON string to dictionary
        logger.info("Successfully extracted JSON from text.")
        return data_dict
    else:
        logger.error("No JSON found.")
        return None


def apply_theme_color(shape, rgb):
    """Apply RGB color to shape fill

    Args:
        shape: The PowerPoint shape to apply color to
        rgb: Tuple of RGB color values (red, green, blue)
    """
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*rgb)

def estimate_text_height(text, font_size, width_inches):
    """Estimate the height needed for text based on font size and width
    Args:
        text (str): The text to estimate height for
        font_size (int): Font size in points
        width_inches (float): Available width in inches

    Returns:
        float: Estimated height in inches needed for the text
    """
    chars_per_line = int((width_inches * 96) / (font_size * 0.6))
    words = text.split()
    lines = 1
    current_line_length = 0
    
    for word in words:
        if current_line_length + len(word) + 1 > chars_per_line:
            lines += 1
            current_line_length = len(word)
        else:
            current_line_length += len(word) + 1
            
    return (lines * font_size * 1.2) / 72

def calculate_content_height(paragraph_text, bullet_points, font_size):
    """Calculate total height needed for content with given font size

    Args:
        paragraph_text (str): Main paragraph text content
        bullet_points (list): List of bullet point text strings
        font_size (int): Font size in points

    Returns:
        float: Total height in inches needed for all content
    """
    total_height = 0
    
    if paragraph_text:
        total_height += estimate_text_height(paragraph_text, font_size, 12)
        total_height += 0.3
    
    if bullet_points:
        for point in bullet_points:
            total_height += estimate_text_height(point, font_size, 11)
            total_height += 0.1
    
    return total_height

def add_bullet_point(tf, text, theme_colors, font_size):
    """Helper function to add properly formatted bullet points
    
    Args:
        tf: Text frame to add bullet point to
        text (str): Text content for the bullet point
        theme_colors (dict): Dictionary of theme colors
        font_size (int): Font size in points
    """
    p = tf.add_paragraph()
    p.text = f"• {text}"
    p.font.size = Pt(font_size)
    
    if any(char in text for char in ["=", "+", "^", "*", "/"]):
        p.font.color.rgb = theme_colors['accent']
    else:
        p.font.color.rgb = theme_colors['text']

    p.space_after = Pt(12)
    p.level = 0
    return p

def create_title_slide(prs, presentation_data, THEME_COLORS):
    """Create title slide with improved text handling

    Args:
        prs: PowerPoint presentation object
        presentation_data (dict): Dictionary containing presentation metadata
        THEME_COLORS (dict): Dictionary of theme colors
    """
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    
    # Background
    background = title_slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    apply_theme_color(background, THEME_COLORS['primary'])
    
    # Accent bar
    accent_bar = title_slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(1), Inches(3),
        Inches(2), Inches(0.1)
    )
    apply_theme_color(accent_bar, THEME_COLORS['accent'])
    
    # Title with improved text handling
    title_box = title_slide.shapes.add_textbox(
        Inches(1), Inches(2),
        Inches(11), Inches(1.5)  # Increased height for text wrapping
    )
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    title_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT  # Enable auto-sizing
    
    title_para = title_frame.add_paragraph()
    title_para.text = presentation_data.get('title', 'Presentation Title')
    title_para.font.size = Pt(54)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(255, 255, 255)
    
    # Calculate if we need to reduce font size for long titles
    while (estimate_text_height(title_para.text, 54, 11) > 1.5 and 
           title_para.font.size > Pt(32)):
        current_size = title_para.font.size.pt
        title_para.font.size = Pt(current_size - 2)
    
    # Subtitle with improved text handling
    subtitle_box = title_slide.shapes.add_textbox(
        Inches(1), Inches(3.5),
        Inches(11), Inches(2)  # Increased height for text wrapping
    )
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.word_wrap = True
    subtitle_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
    
    subtitle_para = subtitle_frame.add_paragraph()
    subtitle_para.text = presentation_data.get('subtitle', '')
    subtitle_para.font.size = Pt(32)
    subtitle_para.font.color.rgb = THEME_COLORS['secondary']
    
    # Calculate if we need to reduce font size for long subtitles
    while (estimate_text_height(subtitle_para.text, 32, 11) > 2 and 
           subtitle_para.font.size > Pt(24)):
        current_size = subtitle_para.font.size.pt
        subtitle_para.font.size = Pt(current_size - 2)

def create_content_slide(prs, slide_data, idx, THEME_COLORS):
    """Create a single content slide with dynamic font sizing
    
    Args:
        prs (Presentation): PowerPoint presentation object
        slide_data (dict): Data for the slide including title, paragraph and bullet points
        idx (int): Slide index number
        THEME_COLORS (dict): Theme color definitions
    """
   
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    # Background
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0,
        prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(248, 248, 248)
    
    # Left bar
    left_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0, 0, Inches(0.25), prs.slide_height
    )
    apply_theme_color(left_bar, THEME_COLORS['primary'])
    
    # Title (keep original size)
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.5),
        Inches(12), Inches(1)
    )
    title_frame = title_box.text_frame
    title_para = title_frame.add_paragraph()
    title_para.text = slide_data.get('title', f'Slide {idx}')
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = THEME_COLORS['primary']
    
    # Accent bar
    accent_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5), Inches(1.4),
        Inches(2), Inches(0.06)
    )
    apply_theme_color(accent_bar, THEME_COLORS['accent'])
    
    # Calculate optimal font size for content
    paragraph_text = slide_data.get('paragraph', '')
    bullet_points = slide_data.get('bullet_points', [])
    
    font_size = 20
    MIN_FONT_SIZE = 14
    MAX_CONTENT_HEIGHT = 4.5
    
    while (font_size > MIN_FONT_SIZE and 
           calculate_content_height(paragraph_text, bullet_points, font_size) > MAX_CONTENT_HEIGHT):
        font_size -= 1
    
    current_y = Inches(1.8)
    
    # Add paragraph with calculated font size
    if paragraph_text:
        content_box = slide.shapes.add_textbox(
            Inches(0.5), current_y,
            Inches(12), Inches(4)
        )
        tf = content_box.text_frame
        tf.word_wrap = True
        
        p = tf.add_paragraph()
        p.text = paragraph_text
        p.font.size = Pt(font_size)
        p.font.color.rgb = THEME_COLORS['text']
        p.space_after = Pt(24)
        
        current_y += Inches(estimate_text_height(paragraph_text, font_size, 12) + 0.3)
    
    # Add bullet points with calculated font size
    if bullet_points:
        bullet_box = slide.shapes.add_textbox(
            Inches(0.5), current_y,
            Inches(12), Inches(4.5)
        )
        tf = bullet_box.text_frame
        tf.word_wrap = True
        
        for point in bullet_points:
            add_bullet_point(tf, point, THEME_COLORS, font_size)
    
    # Slide number
    slide_number = slide.shapes.add_textbox(
        Inches(12), Inches(6.8),
        Inches(0.5), Inches(0.3)
    )
    slide_number_frame = slide_number.text_frame
    slide_number_para = slide_number_frame.add_paragraph()
    slide_number_para.text = str(idx)
    slide_number_para.font.size = Pt(14)
    slide_number_para.font.color.rgb = THEME_COLORS['primary']
    slide_number_para.alignment = PP_ALIGN.RIGHT

def create_presentation(json_data):
    """Create a PowerPoint presentation from JSON data.
    
    Args:
        json_data (dict): JSON data containing presentation content and structure
    
    Returns:
        Presentation: A PowerPoint presentation object
    """
    prs = Presentation()
    
    THEME_COLORS = {
        'primary': RGBColor(0, 75, 135),
        'secondary': RGBColor(240, 240, 240),
        'accent': RGBColor(255, 127, 0),
        'text': RGBColor(51, 51, 51)
    }
    
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    presentation_data = json_data.get('presentation', {})
    
    # Create title slide with improved handling
    create_title_slide(prs, presentation_data, THEME_COLORS)
    
    # Create content slides
    for idx, slide_data in enumerate(presentation_data.get('slides', []), 1):
        create_content_slide(prs, slide_data, idx, THEME_COLORS)
    
    return prs